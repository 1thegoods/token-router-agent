"""Generate a professional PowerPoint for the Token Router Agent project."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
BG_DARK = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x6C, 0x5C, 0xE7)
ACCENT_LIGHT = RGBColor(0xA2, 0x9B, 0xFE)
WHITE = RGBColor(0xE8, 0xE8, 0xF0)
MUTED = RGBColor(0x88, 0x88, 0xA8)
SUCCESS = RGBColor(0x00, 0xCE, 0xC9)
WARNING = RGBColor(0xFD, 0xCB, 0x6E)
ERROR = RGBColor(0xFF, 0x6B, 0x6B)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, icon=">>>"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return txBox


# ====== SLIDE 1 - Title ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_text(slide, Inches(1.5), Inches(2.8), Inches(10.333), Inches(1), "Token Router Agent", font_size=54, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(3.8), Inches(9.333), Inches(0.8), "An Intelligent Multi-Model Routing System for Local AI Agents", font_size=22, color=MUTED, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.667), Inches(4.8), Inches(2), Inches(0.03), ACCENT)
add_text(slide, Inches(2), Inches(5.2), Inches(9.333), Inches(0.5), "By Steven Hanna", font_size=20, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(6.5), Inches(9.333), Inches(0.4), "100% Local  |  Zero API Cost  |  Multi-Model Intelligence", font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)


# ====== SLIDE 2 - The Problem ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "The Problem", font_size=36, bold=True, color=ACCENT_LIGHT)

problems = [
    ("Expensive API Costs", "Cloud AI APIs (OpenAI, Anthropic) charge per token.\nHeavy users can spend $50-$500+/month."),
    ("One-Size-Fits-All Routing", "Most apps send every query to the same large model.\n'Hello' costs the same as 'Write me a compiler.'"),
    ("Privacy Concerns", "Sensitive data gets sent to third-party servers.\nNo control over data retention or usage."),
    ("No Autonomous Execution", "Chatbots just give instructions - they don't\nactually run code or save files for you."),
]
for i, (title, desc) in enumerate(problems):
    x = Inches(0.8 + (i % 2) * 6.2)
    y = Inches(1.5 + (i // 2) * 2.8)
    card = add_shape(slide, x, y, Inches(5.8), Inches(2.4), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40))
    add_text(slide, x + Inches(0.4), y + Inches(0.35), Inches(5), Inches(0.5), title, font_size=20, bold=True)
    add_text(slide, x + Inches(0.4), y + Inches(1.1), Inches(5), Inches(1.2), desc, font_size=14, color=MUTED)


# ====== SLIDE 3 - Our Solution ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Our Solution: Token Router Agent", font_size=36, bold=True, color=ACCENT_LIGHT)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6), "A locally-hosted AI agent that intelligently routes every task to the cheapest capable model - saving tokens while maximizing quality.", font_size=17, color=MUTED)

features = [
    ("Smart Classifier", "Zero-token heuristic classifier analyzes each prompt\nand assigns difficulty: TRIVIAL > EASY > MEDIUM > HARD"),
    ("Intelligent Router", "Routes simple tasks to small models, complex tasks to\na multi-model committee - all running locally via Ollama"),
    ("Autonomous Agent", "Doesn't just chat - it executes. Writes files, runs\ncommands, searches the web, and saves results to your PC"),
    ("Beautiful Web UI", "Full chat interface with session history, file uploads,\ndrag-and-drop, and real-time model/token stats"),
]
for i, (title, desc) in enumerate(features):
    x = Inches(0.8 + (i % 2) * 6.2)
    y = Inches(2.2 + (i // 2) * 2.5)
    card = add_shape(slide, x, y, Inches(5.8), Inches(2.1), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40))
    add_text(slide, x + Inches(0.4), y + Inches(0.3), Inches(4), Inches(0.5), title, font_size=20, bold=True, color=SUCCESS)
    add_text(slide, x + Inches(0.4), y + Inches(1.0), Inches(5), Inches(1), desc, font_size=14, color=MUTED)


# ====== SLIDE 4 - Architecture ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "System Architecture", font_size=36, bold=True, color=ACCENT_LIGHT)

flow_steps = [
    ("User Prompt", ACCENT, Inches(0.5)),
    ("Classifier (Zero-Token)", RGBColor(0x00, 0x80, 0x80), Inches(2.6)),
    ("Router (Brain)", ACCENT, Inches(4.7)),
]
for label, color, y in flow_steps:
    add_shape(slide, Inches(1.5), y, Inches(3.5), Inches(1.6), color)
    add_text(slide, Inches(1.7), y + Inches(0.4), Inches(3.1), Inches(0.8), label, font_size=18, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(2.1), Inches(1.5), Inches(0.4), "V", font_size=24, color=MUTED, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(2.5), Inches(4.2), Inches(1.5), Inches(0.4), "V", font_size=24, color=MUTED, alignment=PP_ALIGN.CENTER)

routes = [
    ("TRIVIAL / EASY", "Single small local model (e.g., Gemma 4)", SUCCESS, Inches(0.8)),
    ("MEDIUM", "2-model committee + synthesizer", WARNING, Inches(2.8)),
    ("HARD", "3-model full committee + synthesizer + escalation", ERROR, Inches(4.8)),
]
for label, desc, color, y in routes:
    add_shape(slide, Inches(6.5), y, Inches(6), Inches(1.5), BG_CARD, border_color=color)
    add_text(slide, Inches(6.8), y + Inches(0.2), Inches(2.5), Inches(0.5), label, font_size=18, bold=True, color=color)
    add_text(slide, Inches(6.8), y + Inches(0.7), Inches(5.4), Inches(0.7), desc, font_size=14, color=MUTED)

add_text(slide, Inches(5.2), Inches(3.0), Inches(1.2), Inches(0.4), "-->", font_size=28, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)


# ====== SLIDE 5 - Committee System ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "Multi-Model Committee System", font_size=36, bold=True, color=ACCENT_LIGHT)
add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6), "For complex tasks, multiple local models collaborate to produce a higher-quality answer than any single model could alone.", font_size=17, color=MUTED)

model_list = [
    ("Gemma 4", "Google's latest efficient model"),
    ("Qwen 3 8B", "Alibaba's powerful reasoning model"),
    ("DeepSeek R1", "Chain-of-thought specialist"),
]
for i, (name, desc) in enumerate(model_list):
    x = Inches(0.8 + i * 3.2)
    add_shape(slide, x, Inches(2.2), Inches(2.8), Inches(2.0), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40))
    add_text(slide, x + Inches(0.3), Inches(2.5), Inches(2.2), Inches(0.5), name, font_size=18, bold=True)
    add_text(slide, x + Inches(0.3), Inches(3.1), Inches(2.2), Inches(0.8), desc, font_size=13, color=MUTED)
    add_text(slide, x + Inches(0.8), Inches(4.2), Inches(1.2), Inches(0.4), "V Draft", font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0.8), Inches(4.8), Inches(9.4), Inches(1.8), ACCENT, border_color=ACCENT_LIGHT)
add_text(slide, Inches(1.2), Inches(5.0), Inches(8.5), Inches(0.5), "Synthesizer Model", font_size=22, bold=True)
add_text(slide, Inches(1.2), Inches(5.6), Inches(8.5), Inches(0.8), "Reads all drafts, identifies the best reasoning from each, and produces a single\nhigh-quality final answer. If confidence is low, triggers escalation.", font_size=15, color=RGBColor(0xD0, 0xD0, 0xE0))
add_text(slide, Inches(1.2), Inches(6.3), Inches(8.5), Inches(0.4), "Confidence Score:  >= 0.6 = Return  |  < 0.6 = Escalate or Retry", font_size=13, color=WARNING)


# ====== SLIDE 6 - Agent Tools ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Autonomous Agent Tools", font_size=36, bold=True, color=ACCENT_LIGHT)
add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6), "Unlike a chatbot, Token Router Agent can execute real actions on your computer:", font_size=17, color=MUTED)

tools = [
    ("read_file", "Read any file on your system"),
    ("write_file", "Create or overwrite files"),
    ("run_command", "Execute terminal commands"),
    ("list_dir", "Browse directory contents"),
    ("search_files", "Find files by pattern"),
    ("search_web", "Search the internet (DuckDuckGo)"),
    ("read_url", "Fetch and read any webpage"),
]
for i, (name, desc) in enumerate(tools):
    x = Inches(0.8)
    y = Inches(2.0 + i * 0.72)
    add_shape(slide, x, y, Inches(11.5), Inches(0.6), BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40))
    add_text(slide, x + Inches(0.5), y + Inches(0.1), Inches(2.5), Inches(0.4), name, font_size=16, bold=True, color=ACCENT_LIGHT, font_name="Consolas")
    add_text(slide, x + Inches(3.5), y + Inches(0.1), Inches(7), Inches(0.4), desc, font_size=15, color=MUTED)


# ====== SLIDE 7 - Tech Stack ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Technology Stack", font_size=36, bold=True, color=ACCENT_LIGHT)

stack = [
    ("Backend", ["Python 3.11 + Flask", "Ollama (local model runtime)", "ThreadPoolExecutor (parallel)"], SUCCESS),
    ("Frontend", ["HTML5 / CSS3 / Vanilla JS", "Glassmorphism UI design", "Drag-and-drop file uploads"], ACCENT_LIGHT),
    ("AI Models", ["Gemma 4 e2b (Google)", "Qwen 3 8B (Alibaba)", "DeepSeek R1 8B"], WARNING),
    ("Tools", ["DuckDuckGo (web search)", "subprocess (commands)", "pathlib (filesystem ops)"], ERROR),
]
for i, (category, items, color) in enumerate(stack):
    x = Inches(0.8 + i * 3.15)
    add_shape(slide, x, Inches(1.5), Inches(2.85), Inches(5.2), BG_CARD, border_color=color)
    add_text(slide, x + Inches(0.3), Inches(1.7), Inches(2.3), Inches(0.5), category, font_size=20, bold=True, color=color)
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.3), Inches(2.5 + j * 0.6), Inches(2.3), Inches(0.5), "- " + item, font_size=13, color=MUTED)


# ====== SLIDE 8 - Comparison Table ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "Why Token Router Agent?", font_size=36, bold=True, color=ACCENT_LIGHT)

comparisons = [
    ("Feature", "ChatGPT / Cloud AI", "Token Router Agent"),
    ("Cost", "$20-$200+/month", "FREE (100% local)"),
    ("Privacy", "Data sent to cloud", "Everything stays on your PC"),
    ("Speed (Simple)", "~2s (network latency)", "~0.5s (local inference)"),
    ("Autonomous Actions", "No (chat only)", "Yes (files, commands, web)"),
    ("Multi-Model", "Single model", "Committee of 3 models"),
    ("Customizable", "Limited", "Full control (open source)"),
    ("Offline Support", "No", "Yes (local models)"),
]
for i, (feat, cloud, ours) in enumerate(comparisons):
    y = Inches(1.5 + i * 0.68)
    bg = BG_CARD if i > 0 else RGBColor(0x20, 0x20, 0x35)
    add_shape(slide, Inches(0.8), y, Inches(3.5), Inches(0.55), bg, border_color=RGBColor(0x2A, 0x2A, 0x40))
    add_shape(slide, Inches(4.4), y, Inches(4), Inches(0.55), bg, border_color=RGBColor(0x2A, 0x2A, 0x40))
    add_shape(slide, Inches(8.5), y, Inches(4), Inches(0.55), bg, border_color=RGBColor(0x2A, 0x2A, 0x40))
    fc = WHITE if i == 0 else MUTED
    add_text(slide, Inches(1.0), y + Inches(0.08), Inches(3.1), Inches(0.4), feat, font_size=14, bold=(i==0), color=fc)
    add_text(slide, Inches(4.6), y + Inches(0.08), Inches(3.6), Inches(0.4), cloud, font_size=14, bold=(i==0), color=ERROR if i > 0 else fc)
    add_text(slide, Inches(8.7), y + Inches(0.08), Inches(3.6), Inches(0.4), ours, font_size=14, bold=(i==0), color=SUCCESS if i > 0 else fc)


# ====== SLIDE 9 - How It Works ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "How It Works - Example", font_size=36, bold=True, color=ACCENT_LIGHT)

demo_steps = [
    ('1', 'User sends prompt', '"Search the web for Apple stock price, create a graph script, save to desktop"', ACCENT),
    ('2', 'Classifier analyzes', 'Type: CODING  |  Difficulty: HARD  |  Estimated tokens: 1200', WARNING),
    ('3', 'Committee activates', '3 local models generate independent drafts. Synthesizer merges into best answer.', SUCCESS),
    ('4', 'Agent executes tools', 'search_web -> write_file -> run_command. Actually saves the script to your Desktop.', ACCENT_LIGHT),
    ('5', 'Response delivered', 'User sees the completed work in the chat with model info, tokens used, and latency.', SUCCESS),
]
for i, (num, title, desc, color) in enumerate(demo_steps):
    y = Inches(1.4 + i * 1.15)
    add_shape(slide, Inches(1.0), y + Inches(0.1), Inches(0.6), Inches(0.6), color)
    add_text(slide, Inches(1.0), y + Inches(0.15), Inches(0.6), Inches(0.5), num, font_size=20, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.0), y + Inches(0.05), Inches(4), Inches(0.4), title, font_size=18, bold=True)
    add_text(slide, Inches(2.0), y + Inches(0.45), Inches(10), Inches(0.7), desc, font_size=13, color=MUTED)


# ====== SLIDE 10 - Thank You ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "What's Next", font_size=36, bold=True, color=ACCENT_LIGHT)

future_items = [
    "Vision model support (analyze images and screenshots)",
    "Streaming responses for real-time output",
    "Plugin system for custom tools (email, calendar, APIs)",
    "Multi-user support with authentication",
    "Model fine-tuning integration for personalized responses",
    "Voice input/output support",
]
add_bullet_list(slide, Inches(1.0), Inches(1.4), Inches(10), Inches(3.5), future_items, font_size=18)

add_text(slide, Inches(1.5), Inches(5.0), Inches(10.333), Inches(0.8), "Thank You!", font_size=48, bold=True, alignment=PP_ALIGN.CENTER, color=ACCENT_LIGHT)
add_text(slide, Inches(2), Inches(5.9), Inches(9.333), Inches(0.5), "Token Router Agent - Built by Steven Hanna", font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(6.5), Inches(9.333), Inches(0.4), "100% Open Source  |  100% Local  |  100% Free", font_size=14, color=SUCCESS, alignment=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.expanduser("~"), "Desktop", "Token_Router_Agent_Presentation.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
